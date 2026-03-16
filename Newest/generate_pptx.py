import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()

# Set slide width and height to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Setup colors
DARK_BLUE = RGBColor(25, 39, 93)   # 19275D
ORANGE = RGBColor(226, 85, 44)     # E2552C
LIGHT_BLUE = RGBColor(130, 170, 221) # 82AADD
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(17, 28, 36)        # 111C24

# Helpers
def add_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_text(slide, text, left, top, width, height, font_size, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Outfit'
    return txBox

def add_image(slide, img_path, left, top, width=None, height=None):
    if os.path.exists(img_path):
        if width and height:
            slide.shapes.add_picture(img_path, left, top, width, height)
        else:
            slide.shapes.add_picture(img_path, left, top)

# -------------------------------------------------------------------------
# Slide 1: Hero
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide1, DARK_BLUE)

try:
    add_image(slide1, "Увеличенный NEWEST логотип.png", Inches(5.16), Inches(2), Inches(3), Inches(1))
except Exception as e:
    print(f"Warning: Logo not found or failed to load. {e}")

add_text(slide1, "NEWEST", Inches(1), Inches(3.5), Inches(11.33), Inches(1.5), 72, WHITE, True, PP_ALIGN.CENTER)
add_text(slide1, "SMART HOME", Inches(1), Inches(4.5), Inches(11.33), Inches(0.5), 24, ORANGE, True, PP_ALIGN.CENTER)
add_text(slide1, "Коммерческое предложение для Мечта Маркет", Inches(1), Inches(6), Inches(11.33), Inches(0.5), 18, LIGHT_BLUE, False, PP_ALIGN.CENTER)

# -------------------------------------------------------------------------
# Slide 2: О компании
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, DARK)

# Image on the right half
img2 = "../Презентация Newest/03909051-8bcb-47bc-8a4c-fdb84da649e2.jpeg"
try:
    # crop or fit image on right half
    add_image(slide2, img2, Inches(7.5), Inches(0), Inches(5.833), Inches(7.5))
except Exception as e:
    print(f"Warning: Image 2 not loaded. {e}")

add_text(slide2, "О КОМПАНИИ", Inches(1), Inches(1), Inches(6), Inches(0.5), 14, ORANGE, True)
add_text(slide2, "Кто мы", Inches(1), Inches(1.5), Inches(6), Inches(0.8), 44, WHITE, True)
desc = "Newest — казахстанский бренд аксессуаров для бытовой техники.\n\nМы специализируемся на защитных чехлах для пультов дистанционного управления ведущих мировых брендов.\n\nНаша миссия — продлить срок службы пультов ТВ и защитить их от повседневного износа, пыли и повреждений."
add_text(slide2, desc, Inches(1), Inches(2.5), Inches(5.5), Inches(3), 16, WHITE)

def add_card(slide, num, text, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    add_text(slide, num, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.5), 24, ORANGE, True)
    add_text(slide, text, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), Inches(0.5), 10, WHITE)

add_card(slide2, "2", "Ведущих бренда Samsung & LG", Inches(1), Inches(5.5), Inches(2.5), Inches(1.2))
add_card(slide2, "50+", "Моделей пультов в каталоге", Inches(3.7), Inches(5.5), Inches(2.5), Inches(1.2))


# -------------------------------------------------------------------------
# Slide 3: Преимущества (Skills)
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, WHITE)
add_text(slide3, "ПРЕИМУЩЕСТВА", Inches(1), Inches(0.8), Inches(6), Inches(0.5), 14, ORANGE, True)
add_text(slide3, "Почему это работает в рознице", Inches(1), Inches(1.3), Inches(8), Inches(0.8), 44, DARK_BLUE, True)

benefits = [
    ("01", "Импульсная покупка", "Цена до 5 000 ₸ — покупатель берет не задумываясь вместе с телевизором."),
    ("02", "Кросс-продажа", "Идеальный допник к каждому ТВ. Конверсия на кассе 15-25%."),
    ("03", "Нулевой возврат", "Простой продукт без брака. Возвраты менее 0.5%."),
    ("04", "Маржа 100%", "РРЦ 5 000 ₸ при закупке 2 500 ₸. Удвоение вложений."),
    ("05", "Компактность", "Минимум места на полке. 100 штук — одна маленькая коробка."),
    ("06", "Отсрочка 30 дней", "Никаких предоплат. Получаете товар, продаете, оплачиваете по факту.")
]

for i, b in enumerate(benefits):
    col = i % 3
    row = i // 3
    left = Inches(1 + col * 3.9)
    top = Inches(2.7 + row * 2.2)
    
    shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.5), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 246, 250)
    shape.line.color.rgb = LIGHT_BLUE
    
    add_text(slide3, b[0], left + Inches(0.2), top + Inches(0.2), Inches(1), Inches(0.4), 20, ORANGE, True)
    add_text(slide3, b[1], left + Inches(0.2), top + Inches(0.6), Inches(3.1), Inches(0.4), 16, DARK_BLUE, True)
    add_text(slide3, b[2], left + Inches(0.2), top + Inches(1.0), Inches(3.1), Inches(0.6), 11, DARK_BLUE)


# -------------------------------------------------------------------------
# Slide 4: Условия
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, DARK_BLUE)

try:
    img4 = "../Презентация Newest/bb7fa294-14b4-4185-95d1-1a549f0d99db.jpeg"
    add_image(slide4, img4, Inches(8.5), Inches(0.8), Inches(4.5), Inches(5.9))
except Exception as e:
    print(f"Warning: Image 4 not loaded. {e}")

add_text(slide4, "КОММЕРЧЕСКОЕ ПРЕДЛОЖЕНИЕ", Inches(1), Inches(0.8), Inches(6), Inches(0.5), 14, LIGHT_BLUE, True)
add_text(slide4, "Условия сотрудничества", Inches(1), Inches(1.3), Inches(6), Inches(0.8), 44, WHITE, True)

# Card 1: Wholesale
shape1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.5), Inches(7), Inches(2.6))
shape1.fill.solid()
shape1.fill.fore_color.rgb = WHITE
shape1.line.fill.background()
add_text(slide4, "ОПТОВАЯ ЦЕНА", Inches(1.3), Inches(2.8), Inches(6.4), Inches(0.4), 14, LIGHT_BLUE, True)
add_text(slide4, "2 500 ₸", Inches(1.3), Inches(3.2), Inches(6.4), Inches(0.8), 44, DARK_BLUE, True)
add_text(slide4, "• Минимальная партия: 50 шт.\n• Бесплатная доставка по РК\n• Маркетинговая поддержка", Inches(1.3), Inches(4.0), Inches(6.4), Inches(1), 16, DARK_BLUE)

# Card 2: Margin
shape2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.4), Inches(3.4), Inches(1.3))
shape2.fill.solid()
shape2.fill.fore_color.rgb = ORANGE
shape2.line.fill.background()
add_text(slide4, "ПОТЕНЦИАЛ", Inches(1.2), Inches(5.6), Inches(3.0), Inches(0.3), 12, WHITE, True)
add_text(slide4, "Ваша маржа: +100%", Inches(1.2), Inches(5.9), Inches(3.0), Inches(0.6), 18, WHITE, True)

# Card 3: Profit
shape3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(5.4), Inches(3.4), Inches(1.3))
shape3.fill.solid()
shape3.fill.fore_color.rgb = DARK
shape3.line.fill.background()
add_text(slide4, "При продаже 200 шт/мес:", Inches(4.8), Inches(5.6), Inches(3.0), Inches(0.3), 12, LIGHT_BLUE, True)
add_text(slide4, "500 000 ₸ чистыми", Inches(4.8), Inches(5.9), Inches(3.0), Inches(0.6), 18, WHITE, True)


# -------------------------------------------------------------------------
# Slide 5: Contact
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, DARK)

try:
    img5 = "../Презентация Newest/a913c337-3e39-45b3-9ff3-f2cfce38f5e1.jpeg"
    add_image(slide5, img5, Inches(0), Inches(0), Inches(5.333), Inches(7.5))
except Exception as e:
    print(f"Warning: Image 5 not loaded. {e}")

add_text(slide5, "Начнём сотрудничество?", Inches(6), Inches(2.5), Inches(6.33), Inches(0.8), 44, WHITE, True, PP_ALIGN.LEFT)
add_text(slide5, "Готовы обсудить условия и отправить тестовую партию.", Inches(6), Inches(3.6), Inches(6.33), Inches(0.5), 18, LIGHT_BLUE, False, PP_ALIGN.LEFT)

def add_contact(slide, label, val, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), top, Inches(6.5), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    add_text(slide, label + ":", Inches(6.3), top + Inches(0.25), Inches(1.5), Inches(0.4), 14, ORANGE, True, PP_ALIGN.LEFT)
    add_text(slide, val, Inches(7.8), top + Inches(0.25), Inches(4.5), Inches(0.4), 16, WHITE, True, PP_ALIGN.LEFT)

add_contact(slide5, "ПОЧТА", "sales@newest.kz", Inches(4.5))
add_contact(slide5, "САЙТ", "newest.kz", Inches(5.5))

prs.save("Newest_Presentation.pptx")
print("Presentation generated successfully at Newest_Presentation.pptx")
