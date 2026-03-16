import asyncio
import os
import time
from pyppeteer import launch
from pptx import Presentation
from pptx.util import Inches, Pt

async def capture_screenshots():
    browser = await launch(executablePath='/usr/bin/google-chrome', args=['--no-sandbox', '--disable-setuid-sandbox'])
    page = await browser.newPage()
    await page.setViewport({'width': 1920, 'height': 1080})
    await page.goto('http://localhost:8080/presentation.html', {'waitUntil': 'networkidle0'})
    
    # Hide all animations and smooth scrolling
    await page.addStyleTag(content="""
        * { animation: none !important; transition: none !important; }
        html { scroll-behavior: auto !important; }
    """)
    
    slide_ids = ['slide-1', 'slide-2', 'slide-3', 'slide-product', 'slide-packaging', 'slide-4', 'slide-5']
    screenshots = []
    
    for i, slide_id in enumerate(slide_ids):
        print(f"Capturing {slide_id}...")
        element = await page.querySelector(f'#{slide_id}')
        if element:
            await page.evaluate(f"document.getElementById('{slide_id}').scrollIntoView();")
            # Wait for any potential layout shifts
            await asyncio.sleep(1)
            path = f'slide_{i+1}.png'
            await element.screenshot({'path': path})
            screenshots.append(path)
        else:
            print(f"Warning: Element #{slide_id} not found.")
            
    await browser.close()
    return screenshots

def create_pptx(screenshots):
    prs = Presentation()
    # Set slide size to 16:9 (10 x 5.625 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for shot in screenshots:
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        # Add the screenshot image as a background/full image
        slide.shapes.add_picture(shot, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
    prs.save('Presentation_NEWEST.pptx')
    print("PPTX saved as Presentation_NEWEST.pptx")
    # Clean up screenshots
    for shot in screenshots:
        os.remove(shot)

async def main():
    screenshots = await capture_screenshots()
    if screenshots:
        create_pptx(screenshots)
    else:
        print("Error: No screenshots were captured.")

if __name__ == "__main__":
    asyncio.get_event_loop().run_until_complete(main())
